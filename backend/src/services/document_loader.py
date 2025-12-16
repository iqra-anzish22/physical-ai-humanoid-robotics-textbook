import os
import tempfile
from pathlib import Path
from typing import List, Optional, Dict, Any
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from io import BytesIO
import hashlib
from datetime import datetime

from src.models.document import Document, ProcessingStatus
from src.utils.logger import logger


class DocumentLoader:
    """Service for loading documents from various sources."""

    SUPPORTED_FORMATS = {'.pdf', '.docx', '.pptx', '.txt', '.md'}

    def __init__(self):
        self.temp_dir = tempfile.gettempdir()

    def load_document_from_file(
        self,
        file_path: str,
        original_filename: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> Document:
        """
        Load a document from a file path.

        Args:
            file_path: Path to the document file
            original_filename: Original filename for reference
            metadata: Additional metadata to attach to the document

        Returns:
            Document object containing the loaded content
        """
        try:
            # Get file extension
            file_ext = Path(file_path).suffix.lower()

            if file_ext not in self.SUPPORTED_FORMATS:
                raise ValueError(f"Unsupported file format: {file_ext}")

            # Read file content
            with open(file_path, 'rb') as f:
                content_bytes = f.read()

            # Extract text content based on file type
            content = self._extract_content(content_bytes, file_ext)

            # Calculate checksum
            checksum = hashlib.md5(content_bytes).hexdigest()

            # Create document object
            document = Document(
                filename=file_path,
                original_filename=original_filename,
                content=content,
                content_type=file_ext,
                size=len(content_bytes),
                checksum=checksum,
                metadata=metadata or {}
            )

            logger.info(f"Successfully loaded document: {original_filename}")
            return document

        except Exception as e:
            logger.error(f"Error loading document {original_filename}: {str(e)}")
            raise

    def load_document_from_bytes(
        self,
        content_bytes: bytes,
        original_filename: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> Document:
        """
        Load a document from bytes content.

        Args:
            content_bytes: Raw bytes of the document
            original_filename: Original filename for reference
            metadata: Additional metadata to attach to the document

        Returns:
            Document object containing the loaded content
        """
        try:
            # Get file extension from filename
            file_ext = Path(original_filename).suffix.lower()

            if file_ext not in self.SUPPORTED_FORMATS:
                raise ValueError(f"Unsupported file format: {file_ext}")

            # Extract text content based on file type
            content = self._extract_content(content_bytes, file_ext)

            # Calculate checksum
            checksum = hashlib.md5(content_bytes).hexdigest()

            # Save to temporary location
            temp_path = os.path.join(self.temp_dir, f"{checksum}{file_ext}")
            with open(temp_path, 'wb') as f:
                f.write(content_bytes)

            # Create document object
            document = Document(
                filename=temp_path,
                original_filename=original_filename,
                content=content,
                content_type=file_ext,
                size=len(content_bytes),
                checksum=checksum,
                metadata=metadata or {}
            )

            logger.info(f"Successfully loaded document from bytes: {original_filename}")
            return document

        except Exception as e:
            logger.error(f"Error loading document from bytes {original_filename}: {str(e)}")
            raise

    def _extract_content(self, content_bytes: bytes, file_extension: str) -> str:
        """
        Extract text content from document based on file extension.

        Args:
            content_bytes: Raw bytes of the document
            file_extension: File extension indicating document type

        Returns:
            Extracted text content
        """
        if file_extension == '.pdf':
            return self._extract_pdf_content(content_bytes)
        elif file_extension == '.docx':
            return self._extract_docx_content(content_bytes)
        elif file_extension == '.pptx':
            return self._extract_pptx_content(content_bytes)
        elif file_extension in ['.txt', '.md']:
            return content_bytes.decode('utf-8')
        else:
            raise ValueError(f"Unsupported file format: {file_extension}")

    def _extract_pdf_content(self, content_bytes: bytes) -> str:
        """Extract content from PDF file."""
        try:
            pdf_reader = PdfReader(BytesIO(content_bytes))
            text_content = []

            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    text = page.extract_text()
                    if text:
                        text_content.append(text)
                except Exception as e:
                    logger.warning(f"Could not extract text from PDF page {page_num}: {str(e)}")

            return "\n".join(text_content)
        except Exception as e:
            logger.error(f"Error extracting PDF content: {str(e)}")
            raise

    def _extract_docx_content(self, content_bytes: bytes) -> str:
        """Extract content from DOCX file."""
        try:
            doc = docx.Document(BytesIO(content_bytes))
            text_content = []

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)

            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_content.append(cell.text)

            return "\n".join(text_content)
        except Exception as e:
            logger.error(f"Error extracting DOCX content: {str(e)}")
            raise

    def _extract_pptx_content(self, content_bytes: bytes) -> str:
        """Extract content from PPTX file."""
        try:
            presentation = Presentation(BytesIO(content_bytes))
            text_content = []

            for slide_num, slide in enumerate(presentation.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)

            return "\n".join(text_content)
        except Exception as e:
            logger.error(f"Error extracting PPTX content: {str(e)}")
            raise

    def validate_document_format(self, file_path: str) -> bool:
        """Validate if the document format is supported."""
        file_ext = Path(file_path).suffix.lower()
        return file_ext in self.SUPPORTED_FORMATS

    def cleanup_temp_files(self, document: Document):
        """Clean up temporary files associated with the document."""
        try:
            if os.path.exists(document.filename) and document.filename.startswith(self.temp_dir):
                os.remove(document.filename)
                logger.info(f"Cleaned up temporary file: {document.filename}")
        except Exception as e:
            logger.error(f"Error cleaning up temporary file {document.filename}: {str(e)}")